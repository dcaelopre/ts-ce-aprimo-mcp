"""Generate MCP Servers Showcase PowerPoint."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / "MCP-Servers-Showcase.pptx"

ACCENT = RGBColor(0x00, 0x5B, 0x96)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, ACCENT)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(8.4), Inches(1.2))
    stf = sub.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(0xCC, 0xE5, 0xF5)
    sp.alignment = PP_ALIGN.LEFT


def add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(8.4), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    note: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = ACCENT

    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.15), Inches(8.8), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)

    if note:
        nb = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(8.6), Inches(0.6))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(12)
        np.font.italic = True
        np.font.color.rgb = MUTED


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    footer: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = ACCENT

    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        cols,
        Inches(0.6),
        Inches(1.3),
        Inches(8.8),
        Inches(0.4 * (len(rows) + 2)),
    )
    table = table_shape.table

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK

    if footer:
        fb = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(8.8), Inches(0.5))
        fp = fb.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(11)
        fp.font.italic = True
        fp.font.color.rgb = MUTED


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "MCP Servers Showcase",
        "Bringing Aprimo DAM & Salsify PIM into AI assistants\nDiana Carla Elopre",
    )

    add_table_slide(
        prs,
        "What We're Showing Today",
        ["Server", "System", "Purpose"],
        [
            [
                "Aprimo DAM MCP",
                "Digital Asset Management",
                "Search assets, browse taxonomy, inspect metadata & rules",
            ],
            [
                "Salsify MCP",
                "Product Information Management",
                "Read products & assets, explore schema, get API guidance",
            ],
        ],
        'Key idea: Ask Claude — "Find released logos in Aprimo" or "Which products use this image?"',
    )

    add_content_slide(
        prs,
        "What Is MCP?",
        [
            "MCP (Model Context Protocol) = open standard for AI clients to call tools safely",
            "User asks in plain language → Claude chooses the right tool → our server calls the REST API",
            "Our servers are thin, typed wrappers around official Aprimo and Salsify APIs",
            "Users connect once in Claude Desktop; credentials stay on their machine",
        ],
    )

    add_content_slide(
        prs,
        "Why We Built These",
        [
            "Problem: Aprimo and Salsify APIs are powerful but verbose — search syntax and field names trip people up",
            "Problem: AI assistants can't help if they can't reach the system",
            "Solution: Purpose-built MCP tools with clear descriptions so the AI picks the right call",
            "Solution: Read-first design — safe for exploration, troubleshooting, and solution design",
            "Solution: Shared Azure deployment — one URL for the team; each person brings their own credentials",
        ],
    )

    add_table_slide(
        prs,
        "Shared Architecture",
        ["Layer", "Choice"],
        [
            ["Language", "TypeScript"],
            ["Protocol", "MCP SDK + Streamable HTTP (POST /mcp)"],
            ["Hosting", "Azure App Service (GitHub Actions on push to main)"],
            ["Client bridge", "npx mcp-remote in Claude Desktop config"],
            ["Validation", "Zod schemas on every tool input"],
            ["Auth model", "Per-request headers — not stored centrally"],
        ],
        "Security: The shared server is stateless. Credentials travel from your Claude config → server → upstream API.",
    )

    add_section_slide(prs, "Server 1: Aprimo DAM MCP")

    add_table_slide(
        prs,
        "Aprimo DAM MCP — Overview",
        ["", ""],
        [
            ["Version", "1.7.1"],
            ["Scope", "{tenant}.dam.aprimo.com — assets, taxonomy, metadata, rules"],
            ["Endpoint", "celopre-aprimo-mcp-dev…eastus-01.azurewebsites.net/mcp"],
            ["Auth", "OAuth client credentials (Environment, Client ID, Client Secret headers)"],
            ["Repo", "github.com/dcaelopre/ts-ce-aprimo-mcp"],
        ],
        "Not in scope: Aprimo Marketing Operations or other non-DAM products.",
    )

    add_table_slide(
        prs,
        "Aprimo DAM — Tools (4)",
        ["Tool", "What it does"],
        [
            ["search_records", "Find assets by keyword, ID, filters, or advanced search expressions"],
            ["search_classifications", "Browse taxonomy — by name, GUID, or parent node"],
            ["search_field_definitions", "Look up metadata field schema (types, validation, labels)"],
            ["search_rules", "Read DAM automation rules — triggers, conditions, actions"],
        ],
        "Smart default: Returns basic fields unless the user explicitly asks for full metadata.",
    )

    add_content_slide(
        prs,
        "Aprimo — Example Prompts",
        [
            '"Search Aprimo for released images matching brand guidelines."',
            '"Get record {guid} and show title, status, and content type only."',
            '"List child classifications under parent ID {guid}."',
            '"What fields are defined for data type OptionList?"',
            '"Show enabled Record rules and their conditions/actions."',
        ],
        "Demo flow: Search → pick a record → ask about classification → look up a field definition.",
    )

    add_section_slide(prs, "Server 2: Salsify PIM MCP")

    add_table_slide(
        prs,
        "Salsify PIM MCP — Overview",
        ["", ""],
        [
            ["Version", "1.1.0"],
            ["Scope", "Read-only at app.salsify.com/api/v1 — products, assets, property schema"],
            ["Endpoint", "celopre-salsify-mcp-dev.azurewebsites.net/mcp"],
            ["Auth", "Bearer token + org ID (X-Salsify-Api-Token, X-Salsify-Org-Id headers)"],
            ["Repo", "github.com/dcaelopre/ce-ts-salsify-mcp"],
        ],
        "Write operations are documented via recommend_salsify_api_routes — not executed by this server.",
    )

    add_table_slide(
        prs,
        "Salsify PIM — Tools (5)",
        ["Tool", "What it does"],
        [
            ["read_salsify_product", "One product, bulk lookup (≤100 IDs), or filter-based search"],
            ["read_salsify_asset", "One asset, bulk lookup, or filter-based search"],
            ["read_salsify_properties", "Property / attribute schema for the org"],
            ["search_salsify_products_by_asset", "Reverse lookup — which products use this digital asset?"],
            ["recommend_salsify_api_routes", "Advisory only — suggests API routes, payloads, curl examples"],
        ],
        "Differentiator: API route recommender guides integration design without executing write calls.",
    )

    add_content_slide(
        prs,
        "Salsify — Example Prompts",
        [
            '"Read Salsify product SKU-12345."',
            '"Search products where Manufacturer equals Acme."',
            '"Which products are linked to digital asset {hash}?"',
            '"What properties are defined in our Salsify org?"',
            '"Recommend the Salsify API route to create a digital asset from a URL."',
        ],
        "Demo flow: Read product → find main image asset → search linked products → ask for update API route.",
    )

    add_table_slide(
        prs,
        "Side-by-Side Comparison",
        ["", "Aprimo DAM", "Salsify PIM"],
        [
            ["Domain", "Digital assets & DAM admin", "Product catalog & syndication prep"],
            ["Primary use", "Asset discovery, taxonomy, rules", "Product/asset lookup, integration planning"],
            ["Tools", "4 (all execute API calls)", "4 read + 1 advisory"],
            ["Unique capability", "DAM automation rules inspection", "Asset→product reverse search + API catalog"],
            ["Out of scope", "Marketing Operations, campaigns", "Writes, syndication runs, exports"],
        ],
        "Together: Cover the two systems most often involved in content ↔ product integration workflows.",
    )

    add_content_slide(
        prs,
        "How Teammates Connect (~5 min setup)",
        [
            "1. Install Claude Desktop + Node.js 18+",
            "2. Edit claude_desktop_config.json (personal file, per Windows/macOS user)",
            "3. Add the MCP block with your own credentials",
            "4. Restart Claude → verify tools appear (hammer icon)",
            "5. Health check: open server root URL in browser → status: running",
            "Setup guides: docs/claude-desktop-setup.md in each repo. Also works in Cursor IDE.",
        ],
    )

    add_content_slide(
        prs,
        "What We Learned",
        [
            "Tool descriptions matter — tuned to prevent over-fetching (e.g. metadata only when asked)",
            "Per-user credentials scale better than one shared service account",
            "Read-only v1 builds trust; advisory tools bridge to write/integration work safely",
            "Same stack twice — second server took a fraction of the time",
            "Official API docs + Cursor rules keep implementations accurate",
        ],
    )

    add_content_slide(
        prs,
        "Roadmap & Discussion",
        [
            "Possible next steps: additional tools, team rollout, Aprimo↔Salsify integration scenarios",
            "Possible next steps: monitoring, rate limiting, and usage metrics on Azure",
            "Who wants early access / credentials setup help?",
            "Which workflows should we prioritize for live demos?",
            "Any other systems that would benefit from the same MCP pattern?",
        ],
        "Suggested timing: ~15–20 minutes with one live demo per server.",
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Created: {path}")
